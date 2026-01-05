"""Multi-modal document processor for PDFs, images, and mixed content."""

import asyncio
import hashlib
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional, Union
from datetime import datetime

from src.ingestion.multimodal.base import (
    MultiModalProcessor,
    MultiModalContent,
    ModalityType,
)
from src.ingestion.multimodal.image import ImageProcessor
from src.ingestion.multimodal.audio import AudioProcessor

logger = logging.getLogger(__name__)


@dataclass
class DocumentPage:
    """Represents a single page from a document."""
    page_number: int
    text_content: str = ""
    images: list[MultiModalContent] = field(default_factory=list)
    tables: list[dict] = field(default_factory=list)
    bbox: Optional[dict] = None  # Page bounding box


@dataclass
class ExtractedTable:
    """Represents an extracted table."""
    page_number: int
    table_index: int
    headers: list[str]
    rows: list[list[str]]
    bbox: Optional[dict] = None
    confidence: float = 0.0


@dataclass
class DocumentStructure:
    """Represents document structure (TOC, sections, etc.)."""
    title: Optional[str] = None
    author: Optional[str] = None
    created_date: Optional[datetime] = None
    modified_date: Optional[datetime] = None
    sections: list[dict] = field(default_factory=list)  # Hierarchical sections
    table_of_contents: list[dict] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)


class PDFProcessor:
    """Process PDF documents with text, images, and tables."""

    def __init__(
        self,
        extract_images: bool = True,
        extract_tables: bool = True,
        ocr_scanned_pages: bool = True,
        image_processor: Optional[ImageProcessor] = None,
        config: Optional[dict] = None,
    ):
        self.extract_images = extract_images
        self.extract_tables = extract_tables
        self.ocr_scanned_pages = ocr_scanned_pages
        self.image_processor = image_processor or ImageProcessor()
        self.config = config or {}
        self._initialized = False

    async def initialize(self) -> None:
        """Initialize PDF processor."""
        if self.extract_images or self.ocr_scanned_pages:
            await self.image_processor.initialize()
        self._initialized = True

    async def process(
        self,
        pdf_content: Union[bytes, str, Path],
        **kwargs,
    ) -> tuple[list[DocumentPage], DocumentStructure]:
        """Process PDF and extract all content."""
        if not self._initialized:
            await self.initialize()

        # Load PDF
        if isinstance(pdf_content, (str, Path)):
            with open(pdf_content, "rb") as f:
                pdf_bytes = f.read()
        else:
            pdf_bytes = pdf_content

        pages = []
        structure = DocumentStructure()

        try:
            import pymupdf as fitz  # PyMuPDF

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")

            # Extract document metadata
            structure.title = doc.metadata.get("title")
            structure.author = doc.metadata.get("author")
            structure.metadata = dict(doc.metadata)

            # Extract TOC
            toc = doc.get_toc()
            structure.table_of_contents = [
                {"level": item[0], "title": item[1], "page": item[2]}
                for item in toc
            ]

            # Process each page
            for page_num in range(len(doc)):
                page = doc[page_num]
                doc_page = await self._process_page(page, page_num + 1)
                pages.append(doc_page)

            doc.close()

        except ImportError:
            logger.warning("PyMuPDF not installed, trying pdfplumber")
            pages, structure = await self._process_with_pdfplumber(pdf_bytes)

        return pages, structure

    async def _process_page(
        self,
        page,
        page_number: int,
    ) -> DocumentPage:
        """Process a single PDF page."""
        import pymupdf as fitz

        doc_page = DocumentPage(page_number=page_number)

        # Extract text
        text = page.get_text("text")
        doc_page.text_content = text

        # Check if page is scanned (mostly images, little text)
        if self.ocr_scanned_pages and len(text.strip()) < 50:
            # Page might be scanned, try OCR
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_bytes = pix.tobytes("png")

            ocr_text = await self.image_processor.extract_text(img_bytes)
            if ocr_text:
                doc_page.text_content = ocr_text

        # Extract images
        if self.extract_images:
            images = []
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                try:
                    base_image = page.parent.extract_image(xref)
                    img_bytes = base_image["image"]

                    # Process image
                    img_content = await self.image_processor.process(
                        img_bytes,
                        extract_text=True,
                        generate_embedding=True,
                        analyze_visual=False,  # Skip heavy analysis
                    )
                    img_content.metadata["page_number"] = page_number
                    img_content.metadata["image_index"] = img_index
                    images.append(img_content)
                except Exception as e:
                    logger.warning(f"Failed to extract image {img_index}: {e}")

            doc_page.images = images

        # Extract tables
        if self.extract_tables:
            tables = await self._extract_tables_from_page(page, page_number)
            doc_page.tables = tables

        return doc_page

    async def _extract_tables_from_page(
        self,
        page,
        page_number: int,
    ) -> list[dict]:
        """Extract tables from a PDF page."""
        tables = []

        try:
            # Try using tabula or camelot
            import tabula

            # Get page as bytes
            import pymupdf as fitz
            pdf_bytes = page.parent.tobytes()

            dfs = tabula.read_pdf(
                io.BytesIO(pdf_bytes),
                pages=page_number,
                multiple_tables=True,
            )

            for i, df in enumerate(dfs):
                tables.append({
                    "table_index": i,
                    "headers": df.columns.tolist(),
                    "rows": df.values.tolist(),
                    "num_rows": len(df),
                    "num_cols": len(df.columns),
                })

        except ImportError:
            # Fallback: basic table detection
            pass
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")

        return tables

    async def _process_with_pdfplumber(
        self,
        pdf_bytes: bytes,
    ) -> tuple[list[DocumentPage], DocumentStructure]:
        """Process PDF using pdfplumber as fallback."""
        try:
            import pdfplumber

            pages = []
            structure = DocumentStructure()

            with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                structure.metadata = pdf.metadata or {}

                for i, page in enumerate(pdf.pages):
                    doc_page = DocumentPage(page_number=i + 1)
                    doc_page.text_content = page.extract_text() or ""

                    # Extract tables
                    if self.extract_tables:
                        for j, table in enumerate(page.extract_tables()):
                            if table:
                                doc_page.tables.append({
                                    "table_index": j,
                                    "headers": table[0] if table else [],
                                    "rows": table[1:] if len(table) > 1 else [],
                                })

                    pages.append(doc_page)

            return pages, structure

        except ImportError:
            logger.error("Neither PyMuPDF nor pdfplumber is installed")
            return [], DocumentStructure()


class MultiModalDocumentProcessor(MultiModalProcessor):
    """Process any document type with multi-modal content."""

    def __init__(
        self,
        config: Optional[dict] = None,
    ):
        super().__init__(config)
        self.image_processor = ImageProcessor(
            ocr_engine=config.get("ocr_engine", "tesseract") if config else "tesseract",
            embedding_model=config.get("image_embedding_model", "clip") if config else "clip",
            vision_provider=config.get("vision_provider", "local") if config else "local",
        )
        self.audio_processor = AudioProcessor(
            transcription_provider=config.get("transcription_provider", "whisper") if config else "whisper",
            embedding_model=config.get("audio_embedding_model", "wav2vec2") if config else "wav2vec2",
        )
        self.pdf_processor = PDFProcessor(
            image_processor=self.image_processor,
        )

    @property
    def supported_modalities(self) -> list[ModalityType]:
        return [
            ModalityType.TEXT,
            ModalityType.IMAGE,
            ModalityType.AUDIO,
            ModalityType.DOCUMENT,
        ]

    @property
    def supported_formats(self) -> list[str]:
        return (
            [".txt", ".md", ".json", ".xml", ".html"]  # Text
            + self.image_processor.supported_formats  # Images
            + self.audio_processor.supported_formats  # Audio
            + [".pdf", ".docx", ".doc", ".pptx", ".ppt"]  # Documents
        )

    async def initialize(self) -> None:
        """Initialize all processors."""
        await asyncio.gather(
            self.image_processor.initialize(),
            self.audio_processor.initialize(),
            self.pdf_processor.initialize(),
        )
        self._initialized = True

    async def process(
        self,
        content: Union[bytes, str, Path],
        modality: Optional[ModalityType] = None,
        **kwargs,
    ) -> MultiModalContent:
        """Process any supported content type."""
        start_time = datetime.now()

        if not self._initialized:
            await self.initialize()

        # Detect modality if not provided
        if modality is None and isinstance(content, (str, Path)):
            modality = self.detect_modality(content)
        elif modality is None:
            # Try to guess from content
            modality = ModalityType.TEXT

        # Route to appropriate processor
        if modality == ModalityType.IMAGE:
            result = await self.image_processor.process(content, **kwargs)

        elif modality == ModalityType.AUDIO:
            result = await self.audio_processor.process(content, **kwargs)

        elif modality == ModalityType.DOCUMENT:
            result = await self._process_document(content, **kwargs)

        else:  # TEXT
            result = await self._process_text(content, **kwargs)

        result.processing_time_ms = (
            datetime.now() - start_time
        ).total_seconds() * 1000

        return result

    async def _process_text(
        self,
        content: Union[bytes, str, Path],
        **kwargs,
    ) -> MultiModalContent:
        """Process text content."""
        if isinstance(content, bytes):
            text = content.decode("utf-8", errors="ignore")
            source_path = None
        elif isinstance(content, (str, Path)):
            if Path(content).exists():
                source_path = str(content)
                with open(content, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            else:
                text = str(content)
                source_path = None
        else:
            text = str(content)
            source_path = None

        return MultiModalContent(
            id=hashlib.sha256(text.encode()).hexdigest()[:16],
            modality=ModalityType.TEXT,
            source_path=source_path,
            text_content=text,
            file_size_bytes=len(text.encode()),
            mime_type="text/plain",
            processed_at=datetime.now(),
        )

    async def _process_document(
        self,
        content: Union[bytes, str, Path],
        **kwargs,
    ) -> MultiModalContent:
        """Process document (PDF, DOCX, etc.)."""
        if isinstance(content, (str, Path)):
            path = Path(content)
            ext = path.suffix.lower()
            with open(content, "rb") as f:
                doc_bytes = f.read()
            source_path = str(content)
        else:
            doc_bytes = content
            ext = ".pdf"  # Assume PDF for bytes
            source_path = None

        result = MultiModalContent(
            id=hashlib.sha256(doc_bytes).hexdigest()[:16],
            modality=ModalityType.DOCUMENT,
            source_path=source_path,
            file_size_bytes=len(doc_bytes),
            processed_at=datetime.now(),
        )
        result.compute_hash(doc_bytes)

        if ext == ".pdf":
            pages, structure = await self.pdf_processor.process(doc_bytes)

            # Combine all text
            all_text = []
            all_images = []

            for page in pages:
                all_text.append(f"[Page {page.page_number}]\n{page.text_content}")
                all_images.extend(page.images)

                # Add table text
                for table in page.tables:
                    table_text = self._format_table_as_text(table)
                    all_text.append(table_text)

            result.text_content = "\n\n".join(all_text)
            result.metadata["pages"] = len(pages)
            result.metadata["structure"] = {
                "title": structure.title,
                "author": structure.author,
                "toc": structure.table_of_contents,
            }
            result.metadata["images"] = [
                {
                    "id": img.id,
                    "page": img.metadata.get("page_number"),
                    "description": img.image_description,
                }
                for img in all_images
            ]
            result.metadata["tables"] = sum(len(p.tables) for p in pages)

        elif ext in [".docx", ".doc"]:
            result = await self._process_docx(doc_bytes, result)

        elif ext in [".pptx", ".ppt"]:
            result = await self._process_pptx(doc_bytes, result)

        return result

    async def _process_docx(
        self,
        doc_bytes: bytes,
        result: MultiModalContent,
    ) -> MultiModalContent:
        """Process DOCX document."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(doc_bytes))

            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract tables
            tables = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    rows.append(cells)
                if rows:
                    tables.append({
                        "headers": rows[0] if rows else [],
                        "rows": rows[1:] if len(rows) > 1 else [],
                    })

            result.text_content = "\n\n".join(paragraphs)
            result.metadata["paragraphs"] = len(paragraphs)
            result.metadata["tables"] = len(tables)
            result.mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        except ImportError:
            logger.warning("python-docx not installed")

        return result

    async def _process_pptx(
        self,
        doc_bytes: bytes,
        result: MultiModalContent,
    ) -> MultiModalContent:
        """Process PPTX presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(doc_bytes))

            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {i + 1}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))

            result.text_content = "\n\n".join(slides_text)
            result.metadata["slides"] = len(prs.slides)
            result.mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        except ImportError:
            logger.warning("python-pptx not installed")

        return result

    def _format_table_as_text(self, table: dict) -> str:
        """Format table as text for embedding."""
        lines = []
        headers = table.get("headers", [])
        rows = table.get("rows", [])

        if headers:
            lines.append(" | ".join(str(h) for h in headers))
            lines.append("-" * len(lines[0]))

        for row in rows:
            lines.append(" | ".join(str(cell) for cell in row))

        return "\n".join(lines)

    async def extract_text(
        self,
        content: Union[bytes, str, Path],
        **kwargs,
    ) -> str:
        """Extract text from any content type."""
        result = await self.process(content, **kwargs)
        return result.text_content

    async def generate_embedding(
        self,
        content: Union[bytes, str, Path],
        **kwargs,
    ) -> list[float]:
        """Generate embedding for content."""
        result = await self.process(content, **kwargs)

        # Return appropriate embedding
        if result.visual_embedding:
            return result.visual_embedding
        elif result.audio_embedding:
            return result.audio_embedding
        elif result.text_embedding:
            return result.text_embedding
        else:
            # Generate text embedding
            from src.ingestion.embeddings import get_embedder
            embedder = await get_embedder()
            return await embedder.embed(result.text_content)
